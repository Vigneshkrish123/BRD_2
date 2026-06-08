"""
file_parser.py — extract plain text from non-transcript file types.

Pure-Python libraries only — no Microsoft Office, no LibreOffice, no COM automation.
This is the correct approach for cloud/Azure deployment where Office is not installed.

Supported formats:
  .pdf   — pypdf  (text-based PDFs; scanned/image PDFs return empty/partial text)
  .docx  — python-docx (already a project dependency)
  .pptx  — python-pptx
  .xlsx  — openpyxl
  .txt   — handled directly in streamlit_app.py (Teams transcript path)
"""

import io
from pathlib import Path
from loguru import logger


def extract_text(filename: str, raw_bytes: bytes) -> str:
    """
    Extract plain text from a supported file.
    Returns a single string ready for the LLM pipeline.
    Raises ValueError for unsupported extensions.
    """
    ext = Path(filename).suffix.lower()

    if ext == ".pdf":
        return _extract_pdf(filename, raw_bytes)
    elif ext == ".docx":
        return _extract_docx(filename, raw_bytes)
    elif ext == ".pptx":
        return _extract_pptx(filename, raw_bytes)
    elif ext in (".xlsx", ".xls"):
        return _extract_excel(filename, raw_bytes)
    else:
        raise ValueError(f"Unsupported file type '{ext}' in '{filename}'")


# ── PDF ───────────────────────────────────────────────────────────────────────

def _extract_pdf(filename: str, raw_bytes: bytes) -> str:
    try:
        import pypdf
    except ImportError:
        raise RuntimeError("pypdf not installed — run: pip install pypdf")

    reader = pypdf.PdfReader(io.BytesIO(raw_bytes))
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(text)

    if not pages:
        logger.warning(f"PDF '{filename}' yielded no extractable text — may be a scanned image PDF")

    result = "\n\n".join(pages)
    logger.info(f"PDF '{filename}' | {len(reader.pages)} pages | {len(result.split()):,} words extracted")
    return result


# ── Word (.docx) ──────────────────────────────────────────────────────────────

def _extract_docx(filename: str, raw_bytes: bytes) -> str:
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("python-docx not installed — run: pip install python-docx")

    doc = Document(io.BytesIO(raw_bytes))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)

    result = "\n\n".join(paragraphs)
    logger.info(f"Word '{filename}' | {len(paragraphs)} blocks | {len(result.split()):,} words extracted")
    return result


# ── PowerPoint (.pptx) ────────────────────────────────────────────────────────

def _extract_pptx(filename: str, raw_bytes: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx not installed — run: pip install python-pptx")

    prs = Presentation(io.BytesIO(raw_bytes))
    slides = []

    for i, slide in enumerate(prs.slides, start=1):
        slide_lines = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    slide_lines.append(text)
        if slide_lines:
            slides.append(f"[Slide {i}]\n" + "\n".join(slide_lines))

    result = "\n\n".join(slides)
    logger.info(f"PPT '{filename}' | {len(prs.slides)} slides | {len(result.split()):,} words extracted")
    return result


# ── Excel (.xlsx / .xls) ──────────────────────────────────────────────────────

def _extract_excel(filename: str, raw_bytes: bytes) -> str:
    try:
        import openpyxl
    except ImportError:
        raise RuntimeError("openpyxl not installed — run: pip install openpyxl")

    wb = openpyxl.load_workbook(io.BytesIO(raw_bytes), read_only=True, data_only=True)
    sheets = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))

    result = "\n\n".join(sheets)
    logger.info(f"Excel '{filename}' | {len(wb.sheetnames)} sheets | {len(result.split()):,} words extracted")
    return result
